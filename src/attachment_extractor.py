"""Extract text and content from email attachments."""

import io
import logging
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
from docx import Document as WordDocument
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation

from .config import settings
from .models import Attachment, AttachmentType

logger = logging.getLogger(__name__)

# Optional OCR support
try:
    import pytesseract

    HAS_TESSERACT = True
except ImportError:
    HAS_TESSERACT = False
    logger.warning("pytesseract not available - OCR will be disabled")


class AttachmentExtractor:
    """Extract text content from various attachment types."""

    def __init__(self, cache_dir: Optional[Path] = None):
        """Initialize the extractor."""
        self.cache_dir = cache_dir or settings.get_attachment_cache_path()
        self.cache_dir.mkdir(parents=True, exist_ok=True)
        self.max_size = settings.max_attachment_size_mb * 1024 * 1024

    def extract_text(self, data: bytes, attachment: Attachment) -> Optional[str]:
        """Extract text from attachment data."""
        if len(data) > self.max_size:
            logger.warning(
                f"Attachment {attachment.filename} exceeds max size "
                f"({len(data)} > {self.max_size})"
            )
            return None

        try:
            match attachment.attachment_type:
                case AttachmentType.PDF:
                    return self._extract_pdf(data)
                case AttachmentType.WORD:
                    return self._extract_word(data)
                case AttachmentType.EXCEL:
                    return self._extract_excel(data)
                case AttachmentType.POWERPOINT:
                    return self._extract_powerpoint(data)
                case AttachmentType.TEXT:
                    return self._extract_text(data)
                case AttachmentType.IMAGE:
                    # Images are handled separately via OCR/vision
                    return None
                case _:
                    logger.debug(f"Unsupported attachment type: {attachment.attachment_type}")
                    return None
        except Exception as e:
            logger.error(f"Failed to extract text from {attachment.filename}: {e}")
            return None

    def _extract_pdf(self, data: bytes) -> str:
        """Extract text from a PDF file."""
        text_parts = []

        with fitz.open(stream=data, filetype="pdf") as doc:
            for page_num, page in enumerate(doc):
                # Extract text
                page_text = page.get_text()
                if page_text.strip():
                    text_parts.append(f"[Page {page_num + 1}]\n{page_text}")

                # Try OCR on images if text extraction yielded little
                if len(page_text.strip()) < 50 and HAS_TESSERACT:
                    # Render page to image for OCR
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    ocr_text = pytesseract.image_to_string(img)
                    if ocr_text.strip():
                        text_parts.append(f"[Page {page_num + 1} OCR]\n{ocr_text}")

        return "\n\n".join(text_parts)

    def _extract_word(self, data: bytes) -> str:
        """Extract text from a Word document."""
        doc = WordDocument(io.BytesIO(data))
        text_parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)

        return "\n".join(text_parts)

    def _extract_excel(self, data: bytes) -> str:
        """Extract text from an Excel file."""
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        text_parts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"[Sheet: {sheet_name}]"]

            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                if any(v.strip() for v in row_values):
                    sheet_text.append(" | ".join(row_values))

            if len(sheet_text) > 1:
                text_parts.append("\n".join(sheet_text))

        wb.close()
        return "\n\n".join(text_parts)

    def _extract_powerpoint(self, data: bytes) -> str:
        """Extract text from a PowerPoint file."""
        prs = Presentation(io.BytesIO(data))
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Extract from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            slide_text.append(row_text)

            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))

        return "\n\n".join(text_parts)

    def _extract_text(self, data: bytes) -> str:
        """Extract text from a plain text file."""
        # Try common encodings
        for encoding in ["utf-8", "latin-1", "cp1252"]:
            try:
                return data.decode(encoding)
            except UnicodeDecodeError:
                continue

        return data.decode("utf-8", errors="replace")

    def extract_image_for_analysis(self, data: bytes, attachment: Attachment) -> Optional[Image.Image]:
        """Load an image for YOLO/vision analysis."""
        try:
            img = Image.open(io.BytesIO(data))
            # Convert to RGB if necessary (for YOLO)
            if img.mode not in ("RGB", "L"):
                img = img.convert("RGB")
            return img
        except Exception as e:
            logger.error(f"Failed to load image {attachment.filename}: {e}")
            return None

    def ocr_image(self, img: Image.Image) -> Optional[str]:
        """Perform OCR on an image."""
        if not HAS_TESSERACT:
            return None

        try:
            text = pytesseract.image_to_string(img)
            return text.strip() if text.strip() else None
        except Exception as e:
            logger.error(f"OCR failed: {e}")
            return None

    def save_attachment(self, data: bytes, attachment: Attachment, email_id: str) -> Path:
        """Save attachment to cache directory."""
        # Create subdirectory for this email
        email_dir = self.cache_dir / email_id[:16]
        email_dir.mkdir(parents=True, exist_ok=True)

        # Save file
        file_path = email_dir / attachment.filename
        file_path.write_bytes(data)

        return file_path
